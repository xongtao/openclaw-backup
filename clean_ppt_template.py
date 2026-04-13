#!/usr/bin/env python3
"""
智能提取PPT模板 - 清理论文内容，保留模板装饰
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re

def is_content_text(text):
    """判断文本是否是论文内容（而非模板占位符）"""
    if not text or len(text.strip()) == 0:
        return False
    
    text = text.strip()
    
    # 模板占位符/常见标题（保留）
    template_patterns = [
        r'^标题$',
        r'^标题[:：]?\s*$',
        r'^副标题$',
        r'^点击此处添加标题$',
        r'^点击此处添加副标题$',
        r'^目录$',
        r'^Contents?$',
        r'^目录[:：]?\s*$',
        r'^致谢$',
        r'^Thank\s*you$',
        r'^Q\s*&?\s*A$',
        r'^问题与讨论$',
        r'^参考文献$',
        r'^Reference',
        r'^\d+$',  # 纯数字（可能是页码）
        r'^第[一二三四五六七八九十\d]+页',
        r'^Page\s*\d+',
    ]
    
    for pattern in template_patterns:
        if re.match(pattern, text, re.IGNORECASE):
            return False
    
    # 论文内容特征（删除）
    content_patterns = [
        r'图\s*\d+',  # 图1、Figure 1
        r'Fig\.?\s*\d+',
        r'Table\s*\d+',
        r'表\s*\d+',
        r'\d{4}\s*年',  # 年份
        r'\d+\.\d+',  # 小数数据
        r'[\d,]+\s*%',  # 百分比
        r'实验',
        r'结果',
        r'结论',
        r'方法',
        r'材料',
        r'数据',
        r'分析',
        r'比较',
        r'显著',
        r'差异',
        r'P\s*<\s*0\.\d+',  # p值
        r'p\s*value',
        r'显著性',
        r'模型',
        r'算法',
        r'训练',
        r'测试',
        r'准确率',
        r'精度',
        r'result',
        r'method',
        r'conclusion',
        r'data',
        r'analysis',
        r'experiment',
        r'performance',
        r'accuracy',
    ]
    
    for pattern in content_patterns:
        if re.search(pattern, text, re.IGNORECASE):
            return True
    
    # 长文本大概率是内容
    if len(text) > 100:
        return True
    
    # 包含具体数据、句子结构
    if '。' in text or '，' in text or '.' in text:
        return True
    
    return False

def is_likely_content_image(shape, slide_idx):
    """判断图片是否可能是论文内容图（而非背景/logo）"""
    # 获取图片位置和大小
    try:
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
        
        # 小图标/logo类图片（保留）
        if width < Pt(100) and height < Pt(100):
            return False
        
        # 全页背景图（保留）
        slide_width = Pt(960)  # 默认16:9宽度
        slide_height = Pt(540)
        
        if width > Pt(700) and height > Pt(400):
            # 可能是背景
            return False
        
        # 页脚/角落的小图（保留）
        if top > Pt(450) or left > Pt(800):
            if width < Pt(150):
                return False
        
    except:
        pass
    
    # 默认：大尺寸图片可能是内容图
    return True

def clean_ppt_template(src_path, dst_path):
    """清理PPT中的论文内容，保留模板"""
    
    prs = Presentation(src_path)
    
    print(f"原PPT共 {len(prs.slides)} 页")
    
    for idx, slide in enumerate(prs.slides):
        print(f"处理第 {idx+1}/{len(prs.slides)} 页...")
        
        shapes_to_remove = []
        shapes_to_clear_text = []
        
        for shape in slide.shapes:
            # 1. 处理图片 - 判断是否是内容图
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if is_likely_content_image(shape, idx):
                    shapes_to_remove.append(shape)
                    print(f"  删除图片: {shape.name}")
            
            # 2. 处理文本
            elif shape.has_text_frame:
                text_frame = shape.text_frame
                full_text = text_frame.text.strip()
                
                # 检查是否是占位符（保留）
                if shape.is_placeholder:
                    placeholder_type = None
                    try:
                        placeholder_type = shape.placeholder_format.type
                    except:
                        pass
                    
                    # 标题/副标题占位符：清空内容保留样式
                    if placeholder_type in [1, 2, 3]:  # Title, Body, Center Title
                        if is_content_text(full_text):
                            shapes_to_clear_text.append(shape)
                            print(f"  清空占位符文本: {full_text[:30]}...")
                    continue
                
                # 非占位符文本框
                if is_content_text(full_text):
                    # 检查是否是页脚/页码
                    is_footer = False
                    try:
                        if shape.top > Pt(480):  # 靠近底部
                            if re.match(r'^\d+$', full_text) or len(full_text) < 10:
                                is_footer = True
                    except:
                        pass
                    
                    if is_footer:
                        continue  # 保留页脚
                    
                    # 其他内容文本：清空
                    shapes_to_clear_text.append(shape)
                    print(f"  清空文本: {full_text[:40]}...")
        
        # 执行删除
        for shape in shapes_to_remove:
            try:
                sp = shape._element
                sp.getparent().remove(sp)
            except Exception as e:
                print(f"  删除失败: {e}")
        
        # 执行文本清空
        for shape in shapes_to_clear_text:
            try:
                text_frame = shape.text_frame
                # 保留段落结构，清空文本
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = ""
            except Exception as e:
                print(f"  清空失败: {e}")
    
    # 保存
    prs.save(dst_path)
    print(f"\n✅ 模板已保存: {dst_path}")

if __name__ == "__main__":
    src_file = "/data/20251223.pptx"
    dst_file = "/data/模板_清理论文内容.pptx"
    
    print("=" * 60)
    print("🎨 智能PPT模板提取")
    print("   清理论文内容，保留模板装饰")
    print("=" * 60)
    
    clean_ppt_template(src_file, dst_file)
    
    print("\n✨ 完成！")
    print("\n保留项：")
    print("  • 背景图片/配色")
    print("  • 装饰性元素")
    print("  • 页眉/页脚/页码")
    print("  • Logo/学校标识")
    print("  • 标题占位符样式")
    print("\n已清除：")
    print("  • 论文相关的图表")
    print("  • 数据/实验结果")
    print("  • 分析文字内容")
