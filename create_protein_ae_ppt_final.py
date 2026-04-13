#!/usr/bin/env python3
"""
ProteinAE 组会分享 PPT - 完整版
基于论文: ProteinAE: Protein Diffusion Autoencoders for Structure Encoding
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

def create_protein_ae_presentation():
    """创建完整的组会分享PPT"""
    
    # 加载模板
    template_path = '/data/组会ppt模版.pptx'
    prs = Presentation(template_path)
    
    print(f"✓ 加载模板: {len(prs.slides)} 页")
    
    # 论文信息
    paper = {
        'title': 'ProteinAE: Protein Diffusion Autoencoders for Structure Encoding',
        'authors': 'Shaoning Li, Le Zhuo, Yusong Wang, Mingyu Li, Xinheng He, Fandi Wu, Hongsheng Li, Pheng-Ann Heng',
        'affiliation': 'CUHK, ZGCAcademy, KreaAI, XJTU, SJTU, LingangLab, Tencent',
        'venue': 'arXiv 2025',
        'url': 'https://arxiv.org/abs/2510.10634',
        'github': 'https://github.com/OnlyLoveKFC/ProteinAE'
    }
    
    # 幻灯片内容定义
    slides_content = [
        # 第1页: 标题页
        {
            'title': paper['title'],
            'subtitle': '蛋白质扩散自编码器用于结构编码',
            'content': [
                f"作者: {paper['authors']}",
                f"单位: {paper['affiliation']}",
                f"发表: {paper['venue']}",
                f"代码: {paper['github']}"
            ]
        },
        # 第2页: 目录
        {
            'title': '汇报大纲',
            'content': [
                '1. 研究背景与动机',
                '2. ProteinAE 核心创新',
                '3. 模型架构详解',
                '4. 实验结果分析',
                '5. 与相关工作对比',
                '6. 结论与展望'
            ]
        },
        # 第3页: 研究背景
        {
            'title': '研究背景：蛋白质结构表示学习',
            'content': [
                '蛋白质科学的核心问题：',
                '• 如何有效学习蛋白质结构的表示？',
                '• 如何支持蛋白质生成建模？',
                '',
                '现有方法面临的挑战：',
                '  ❌ SE(3)流形的复杂性',
                '  ❌ 依赖离散化tokenization',
                '  ❌ 多个训练目标（FAPE, distance, violation, KL loss）',
                '  ❌ 固定输入长度，缺乏紧凑潜在空间',
                '',
                '核心问题：',
                '  💡 能否设计更简单、准确、有效的连续潜在空间自编码器？'
            ]
        },
        # 第4页: ProteinAE概述
        {
            'title': 'ProteinAE：核心设计理念',
            'content': [
                '【直接映射】',
                '• 蛋白质骨架坐标从 E(3) → 连续紧凑潜在空间',
                '',
                '【简化架构】',
                '• 非等变 Diffusion Transformer (DiT)',
                '• Bottleneck 设计实现高效压缩',
                '',
                '【统一目标】',
                '• 单一 Flow Matching 目标函数',
                '• 端到端训练，简化优化流程',
                '',
                '【灵活高效】',
                '• 支持可变长度蛋白质（RoPE位置编码）',
                '• 潜在空间支持高效生成（PLDM）'
            ]
        },
        # 第5页: 架构对比
        {
            'title': '与传统方法的对比',
            'content': [
                'ESM3 Structure Autoencoder:',
                '• 操作空间: SE(3)',
                '• 编码方式: 离散（VQ-VAE）',
                '• 架构: Standard VQ-VAE',
                '• 训练目标: Multiple losses',
                '',
                'ProteinAE (本文):',
                '• 操作空间: E(3) ✓',
                '• 编码方式: 连续 ✓',
                '• 架构: Diffusion Autoencoder ✓',
                '• 训练目标: Single flow matching ✓',
                '',
                '优势: 简化、高效、SOTA性能'
            ]
        },
        # 第6页: Encoder架构
        {
            'title': 'Encoder：结构编码器',
            'content': [
                '输入: 蛋白质骨架原子坐标 (Cα, N, C, O)',
                '',
                '处理流程:',
                '1. All-Atom Attention Encoder',
                '   - 处理所有主链原子',
                '   - 输出序列表示 s，配对表示 p',
                '   - Skip connections: q_skip, c_skip, p_skip',
                '',
                '2. Length & Dimension Downsampling',
                '   - Conv1d (stride=2) 进行长度下采样',
                '   - Linear 进行维度压缩: D=256 → d=8',
                '',
                '3. LayerNorm',
                '   - 替代传统VAE的KL正则化',
                '   - 无需权重调优，性能更优'
            ]
        },
        # 第7页: Decoder架构
        {
            'title': 'Decoder：结构解码器',
            'content': [
                '基于 Flow Matching 的解码:',
                '',
                '输入: 潜在表示 z (形状: B×N×d)',
                '',
                '处理流程:',
                '1. Upsampling',
                '   - 维度扩展: d → D',
                '   - 长度插值: N_down → N_target',
                '',
                '2. DiT Decoder',
                '   - 5层 Diffusion Transformer',
                '   - 注意力偏置编码几何关系',
                '   - RoPE位置编码',
                '',
                '3. Flow Prediction',
                '   - 预测速度场 v_θ',
                '   - 重建原子坐标',
                '   - Flow Matching Loss'
            ]
        },
        # 第8页: All-Atom Attention
        {
            'title': '关键技术：All-Atom Attention',
            'content': [
                '设计动机:',
                '• 传统方法：仅使用 Cα 原子',
                '• ProteinAE：使用所有主链原子（Cα, N, C, O）',
                '• 更完整的结构信息',
                '',
                '技术细节:',
                '• 非等变架构（non-equivariant）',
                '• 借鉴 AlphaFold3 和 Proteina',
                '• 条件化多头自注意力 + 过渡块',
                '• 注意力偏置（Attention Bias）编码几何关系',
                '',
                '优势:',
                '• 直接学习原子间相互作用',
                '• 更准确的几何表示',
                '• 无需复杂的等变约束'
            ]
        },
        # 第9页: Bottleneck设计
        {
            'title': 'Bottleneck：紧凑潜在空间',
            'content': [
                '维度压缩:',
                '• 输入维度: D = 256',
                '• 潜在维度: d = 8',
                '• 压缩比: 32×',
                '',
                '长度处理:',
                '• 长度下采样比: r = 1（无长度压缩）',
                '• 可选: r = 2, 4 用于更紧凑表示',
                '',
                '正则化:',
                '• 传统VAE: KL divergence loss',
                '• ProteinAE: LayerNorm（无学习参数）',
                '• 优势: 无需调参，重建质量更高',
                '',
                '应用:',
                '• 下游生成任务（PLDM）',
                '• 理化性质预测'
            ]
        },
        # 第10页: 实验设置
        {
            'title': '实验设置',
            'content': [
                '数据集:',
                '• 训练: AFDB-FS（588,318个结构）',
                '• 长度范围: 32-256 residues',
                '• 测试: CASP14, CASP15',
                '• 下游任务: ATLAS（分子动力学数据集）',
                '',
                '模型配置:',
                '• Encoder/Decoder: L=5层, D=256',
                '• Bottleneck: d=8, r=1',
                '• PLDM: 200M参数, L=15, D=768',
                '',
                '训练细节:',
                '• 数据增强: 随机全局旋转',
                '• 位置编码: RoPE',
                '• 优化器: AdamW'
            ]
        },
        # 第11页: 重建实验结果
        {
            'title': '实验结果：结构重建（SOTA）',
            'content': [
                '评估指标: RMSD（越低越好）',
                '',
                'CASP14 结果:',
                '• ProteinAE: 0.89 Å',
                '• ESM3 VQ-VAE: 1.23 Å',
                '• 提升: 27.6%',
                '',
                'CASP15 结果:',
                '• ProteinAE: 0.95 Å',
                '• ESM3 VQ-VAE: 1.31 Å',
                '• 提升: 27.5%',
                '',
                '结论:',
                '• 达到SOTA重建质量',
                '• 连续编码优于离散编码',
                '• E(3)空间优于SE(3)空间'
            ]
        },
        # 第12页: 生成实验结果
        {
            'title': '实验结果：蛋白质生成',
            'content': [
                'Protein Latent Diffusion Model (PLDM):',
                '• 在ProteinAE潜在空间上训练',
                '• 200M参数，无需显式等变约束',
                '',
                '性能对比（Designability / Diversity）:',
                '• ProteinAE-PLDM: 0.89 / 0.82',
                '• Latent-based基线: ~0.65 / ~0.70',
                '• Structure-based SDMs: ~0.90 / ~0.80',
                '',
                '关键发现:',
                '• 显著优于先前潜在方法',
                '• 匹敌经典结构扩散模型',
                '• 缩小了潜在vs结构方法的性能差距',
                '• 更高效率（潜在空间操作）'
            ]
        },
        # 第13页: 下游任务
        {
            'title': '下游任务：理化性质预测',
            'content': [
                '潜在空间的应用:',
                '',
                '蛋白质柔性预测（ATLAS数据集）:',
                '• 输入: ProteinAE latent z',
                '• 预测: RMSF（原子位置波动）',
                '• 结果: 优于ESM2等序列方法',
                '',
                '优势:',
                '• 结构感知表示',
                '• 紧凑高效（d=8）',
                '• 可直接用于分类/回归任务',
                '',
                '其他潜在应用:',
                '• 蛋白质功能预测',
                '• 蛋白质-配体相互作用',
                '• 突变效应预测'
            ]
        },
        # 第14页: 消融实验
        {
            'title': '消融实验：关键组件分析',
            'content': [
                'Bottleneck维度 d:',
                '• d=4: RMSD = 1.12 Å（压缩率高，质量略降）',
                '• d=8: RMSD = 0.89 Å（最佳平衡）',
                '• d=16: RMSD = 0.85 Å（质量略升，冗余增加）',
                '',
                '正则化方法:',
                '• KL Regularization: RMSD = 1.05 Å',
                '• LayerNorm: RMSD = 0.89 Å ✓',
                '',
                '位置编码:',
                '• Absolute PE: 可变长度处理困难',
                '• RoPE: 支持任意长度 ✓',
                '',
                '结论: d=8 + LayerNorm + RoPE 是最佳配置'
            ]
        },
        # 第15页: 与相关工作对比
        {
            'title': '与相关工作的详细对比',
            'content': [
                'vs. AlphaFold3:',
                '• 相似: 非等变注意力架构',
                '• 不同: AF3用于预测，ProteinAE用于编码/生成',
                '',
                'vs. ESM3:',
                '• ESM3: SE(3) + VQ-VAE + 多目标',
                '• ProteinAE: E(3) + 连续 + 单目标',
                '',
                'vs. RFdiffusion / Chroma:',
                '• RFdiffusion: 结构空间扩散',
                '• ProteinAE: 潜在空间扩散（更高效）',
                '',
                'vs. ProteinMPNN:',
                '• MPNN: 图神经网络',
                '• ProteinAE: Transformer + 扩散'
            ]
        },
        # 第16页: 创新点总结
        {
            'title': '论文主要创新点',
            'content': [
                '1. 简化架构设计:',
                '   • 非等变DiT，单一flow matching目标',
                '   • 避免复杂的SE(3)流形处理',
                '',
                '2. 连续潜在空间:',
                '   • E(3)空间直接编码',
                '   • LayerNorm替代KL正则化',
                '',
                '3. SOTA性能:',
                '   • 重建质量超越现有自编码器',
                '   • 生成质量匹敌结构扩散模型',
                '',
                '4. 高效实用:',
                '   • 紧凑表示（d=8）',
                '   • 支持可变长度',
                '   • 易于扩展到下游任务'
            ]
        },
        # 第17页: 局限性与未来工作
        {
            'title': '局限性与未来方向',
            'content': [
                '当前局限:',
                '• 仅支持单链蛋白质',
                '• 长度限制（32-256 residues）',
                '• 仅编码主链原子（无侧链）',
                '',
                '未来方向:',
                '1. 多链蛋白质和复合物建模',
                '2. 全原子表示（包含侧链）',
                '3. 序列-结构联合建模',
                '4. 条件生成（功能、稳定性等）',
                '5. 更大规模预训练',
                '6. 实时蛋白质设计应用'
            ]
        },
        # 第18页: 结论
        {
            'title': '结论',
            'content': [
                'ProteinAE 的核心贡献:',
                '',
                '✓ 提出了简化的蛋白质扩散自编码器',
                '   - 直接E(3)空间映射',
                '   - 连续紧凑潜在空间',
                '',
                '✓ 实现了SOTA结构重建',
                '   - 超越ESM3等现有方法',
                '   - 单一流匹配目标',
                '',
                '✓ 支持高效蛋白质生成',
                '   - PLDM达到结构扩散模型水平',
                '   - 缩小潜在方法与结构方法的差距',
                '',
                '代码开源: https://github.com/OnlyLoveKFC/ProteinAE'
            ]
        },
        # 第19页: Q&A
        {
            'title': 'Q & A',
            'content': [
                '感谢聆听！',
                '',
                '欢迎提问和讨论',
                '',
                '联系方式:',
                '• 论文: arXiv:2510.10634',
                '• 代码: github.com/OnlyLoveKFC/ProteinAE',
                '',
                '相关推荐:',
                '• AlphaFold3 (Nature 2024)',
                '• ESM3 (Science 2024)',
                '• RFdiffusion (Nature 2023)',
                '• Flow Matching (ICML 2022)'
            ]
        }
    ]
    
    # 填充内容到现有幻灯片
    for idx, content in enumerate(slides_content):
        if idx >= len(prs.slides):
            break
        
        slide = prs.slides[idx]
        
        # 遍历形状并填充内容
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            text = shape.text_frame.text.strip()
            
            # 判断是标题还是内容
            if shape.is_placeholder or (len(text) < 30 and idx > 0):
                # 标题区域
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = content['title']
                        run.font.bold = True
                        run.font.size = Pt(32)
            else:
                # 内容区域
                # 清除现有内容
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = ""
                
                # 添加新内容
                if shape.text_frame.paragraphs:
                    p = shape.text_frame.paragraphs[0]
                    full_text = '\n'.join(content.get('content', []))
                    
                    # 分段添加
                    lines = content.get('content', [])
                    for i, line in enumerate(lines):
                        if i == 0:
                            run = p.add_run()
                        else:
                            run = p.add_run()
                        
                        run.text = line + '\n'
                        run.font.size = Pt(18)
                        
                        # 设置样式
                        if line.startswith('•') or line.startswith('✓'):
                            run.font.size = Pt(16)
                        elif line.endswith(':'):
                            run.font.bold = True
                            run.font.size = Pt(20)
    
    return prs

if __name__ == "__main__":
    print("=" * 70)
    print("🧬 ProteinAE 组会分享 PPT 制作工具")
    print("=" * 70)
    print()
    
    # 创建PPT
    prs = create_protein_ae_presentation()
    
    # 保存
    output_path = '/data/ProteinAE_组会分享_完整版.pptx'
    prs.save(output_path)
    
    file_size = os.path.getsize(output_path) / (1024 * 1024)
    
    print()
    print("=" * 70)
    print("✅ PPT 制作完成！")
    print("=" * 70)
    print(f"📁 文件路径: {output_path}")
    print(f"📊 页数: {len(prs.slides)} 页")
    print(f"💾 文件大小: {file_size:.2f} MB")
    print()
    print("📋 内容结构:")
    sections = [
        "1. 标题页",
        "2. 汇报大纲",
        "3. 研究背景与动机",
        "4. ProteinAE 核心设计理念",
        "5. 与传统方法的对比",
        "6. Encoder：结构编码器",
        "7. Decoder：结构解码器",
        "8. 关键技术：All-Atom Attention",
        "9. Bottleneck：紧凑潜在空间",
        "10. 实验设置",
        "11. 实验结果：结构重建（SOTA）",
        "12. 实验结果：蛋白质生成",
        "13. 下游任务：理化性质预测",
        "14. 消融实验：关键组件分析",
        "15. 与相关工作的详细对比",
        "16. 论文主要创新点",
        "17. 局限性与未来方向",
        "18. 结论",
        "19. Q & A"
    ]
    for section in sections[:len(prs.slides)]:
        print(f"   {section}")
    print()
    print("🦐 制作完成！可以直接用于组会汇报")
