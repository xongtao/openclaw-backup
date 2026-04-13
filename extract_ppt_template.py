#!/usr/bin/env python3
"""
提取PPT模板 - 保留母版样式，清理内容
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
import copy
import os

def extract_template(src_path, dst_path):
    """提取PPT模板，清理内容保留样式"""
    
    prs = Presentation(src_path)
    
    # 创建新PPT，使用相同的模板/母版
    new_prs = Presentation(src_path)
    
    # 获取幻灯片数量
    slide_count = len(new_prs.slides)
    print(f"原PPT共 {slide_count} 页幻灯片")
    
    # 我们需要保留母版，但清理幻灯片内容
    # 策略：遍历所有幻灯片，将内容替换为占位符
    
    for idx, slide in enumerate(new_prs.slides):
        print(f"处理第 {idx+1}/{slide_count} 页...")
        
        # 收集需要删除的形状（非母版背景的图片/文本）
        shapes_to_remove = []
        
        for shape in slide.shapes:
            # 跳过母版背景相关的内容
            if shape.is_placeholder:
                continue
                
            # 检查形状类型
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # 图片 - 删除（组会内容相关的图片）
                shapes_to_remove.append(shape)
                
            elif shape.has_text_frame:
                # 文本框 - 清空内容但保留样式
                text_frame = shape.text_frame
                
                # 检查是否是标题或主要文本
                is_title = False
                try:
                    if shape.placeholder_format:
                        ph_type = shape.placeholder_format.type
                        # 1=Title, 2=Body, 3=Center Title, etc.
                        if ph_type in [1, 2, 3]:
                            is_title = True
                except:
                    pass
                
                # 对于非占位符的文本框，替换为占位文本
                if not shape.is_placeholder:
                    # 清空文本但保留样式
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            # 保留格式，替换文本
                            run.text = "[在此输入内容]"
                        if len(paragraph.runs) == 0:
                            # 没有runs，添加一个
                            p = paragraph._p
                            r = p.add_r()
                            rPr = r.get_or_add_rPr()
                            r.text = "[在此输入内容]"
                            
        # 删除图片等非模板元素
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)
    
    # 保存提取的模板
    new_prs.save(dst_path)
    print(f"\n✅ 模板已保存至: {dst_path}")
    return True

def create_clean_template(src_path, dst_path):
    """创建干净的模板 - 只保留母版，清空所有幻灯片内容"""
    
    prs = Presentation(src_path)
    
    # 获取母版信息
    slide_layouts = prs.slide_layouts
    print(f"检测到 {len(slide_layouts)} 种幻灯片布局")
    
    # 创建新的演示文稿
    from pptx import Presentation as NewPresentation
    new_prs = NewPresentation()
    
    # 复制幻灯片宽度/高度设置
    new_prs.slide_width = prs.slide_width
    new_prs.slide_height = prs.slide_height
    
    # 这里有个限制：python-pptx不能直接复制母版
    # 变通方案：保留第一页作为模板示例，清空内容
    
    # 获取第一张幻灯片
    if len(prs.slides) > 0:
        first_slide = prs.slides[0]
        
        # 清空第一页的内容，保留为背景示例
        shapes_to_remove = []
        for shape in first_slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shapes_to_remove.append(shape)
            elif shape.has_text_frame and not shape.is_placeholder:
                # 替换文本
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = ""
        
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)
    
    # 删除其他幻灯片（从后往前删）
    slide_ids = list(range(len(prs.slides) - 1, 0, -1))
    for idx in slide_ids:
        rId = prs.slides._sldIdLst[idx].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[idx]
    
    # 保存
    prs.save(dst_path)
    print(f"\n✅ 干净模板已保存至: {dst_path}")
    print(f"📊 保留了 {len(prs.slides)} 页母版示例")
    return True

if __name__ == "__main__":
    src_file = "/data/20251223.pptx"
    template_file = "/data/模板_保留母版.pptx"
    clean_template = "/data/模板_干净版.pptx"
    
    print("=" * 50)
    print("🎨 PPT 模板提取工具")
    print("=" * 50)
    
    # 方式1：保留所有页面结构，替换内容为占位符
    print("\n【方式1】提取带占位符的模板...")
    extract_template(src_file, template_file)
    
    # 方式2：只保留母版，清空所有内容
    print("\n【方式2】创建干净模板...")
    create_clean_template(src_file, clean_template)
    
    print("\n" + "=" * 50)
    print("✨ 模板提取完成！")
    print(f"\n📁 输出文件：")
    print(f"  1. {template_file}")
    print(f"     - 保留原PPT结构，内容替换为占位符")
    print(f"  2. {clean_template}")
    print(f"     - 仅保留母版和第一页作为示例")
    print("=" * 50)
