#!/usr/bin/env python3
"""
基于ProteinAE论文制作组会分享PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

def create_presentation():
    # 加载模板
    template_path = '/data/组会ppt模版.pptx'
    prs = Presentation(template_path)
    
    # 论文核心内容
    paper_info = {
        'title': 'ProteinAE: Protein Diffusion Autoencoders for Structure Encoding',
        'authors': 'Shaoning Li, Le Zhuo, Yusong Wang, et al.',
        'affiliation': 'CUHK, ZGCAcademy, KreaAI, XJTU, SJTU, LingangLab, Tencent',
        'venue': 'arXiv 2025',
        'github': 'https://github.com/OnlyLoveKFC/ProteinAE'
    }
    
    # 清空现有幻灯片内容（保留布局）
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and not shape.is_placeholder:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = ""
    
    print(f"使用模板创建PPT，共 {len(prs.slides)} 页")
    
    # 第1页: 标题页
    if len(prs.slides) > 0:
        slide = prs.slides[0]
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.lower()
                if 'title' in text or '标题' in text or shape.is_placeholder:
                    # 设置标题
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip():
                                run.text = paper_info['title']
                                break
    
    # 返回处理后的演示文稿
    return prs, paper_info

def add_content_slides(prs, paper_info):
    """添加内容页"""
    
    # 准备内容大纲
    outline = [
        {
            'title': '研究背景与动机',
            'content': [
                '• 蛋白质结构表示学习对蛋白质生成建模至关重要',
                '• 现有方法面临的挑战：',
                '  - SE(3)流形的复杂性',
                '  - 依赖离散化tokenization',
                '  - 需要多个训练目标（FAPE loss, distance loss, violation loss, KL loss等）',
                '  - 固定输入长度，缺乏紧凑的潜在空间',
                '• 核心问题：能否设计一个更简单、准确、有效的连续潜在空间蛋白质自编码器？'
            ]
        },
        {
            'title': 'ProteinAE 核心创新',
            'content': [
                '• 直接在 E(3) 空间操作，避免离散化',
                '• 非等变扩散 Transformer 架构（DiT）',
                '• 单一训练目标：Flow Matching',
                '• 紧凑的连续潜在空间（bottleneck设计）',
                '• 端到端训练，简化优化流程',
                '• 支持可变长度蛋白质输入（RoPE位置编码）'
            ]
        },
        {
            'title': '模型架构',
            'content': [
                'Encoder（编码器）：',
                '• All-Atom Attention 模块处理输入',
                '• DiT Stack 进行特征提取',
                '• Length & Dimension Downsampling',
                '• LayerNorm 替代 KL 正则化',
                '',
                'Decoder（解码器）：',
                '• Latent 上采样和扩展',
                '• Flow-based 结构重建',
                '• 预测速度场 v_θ 进行去噪'
            ]
        },
        {
            'title': '关键技术细节',
            'content': [
                '1. 非等变架构设计：',
                '   - 借鉴 AlphaFold3 和 Proteina',
                '   - 使用注意力偏置（attention bias）编码几何关系',
                '',
                '2. All-Atom Attention：',
                '   - 处理所有主链原子（Cα, N, C, O）',
                '   - 输出序列表示和配对表示',
                '',
                '3. Bottleneck 设计：',
                '   - 长度下采样比例 r=1（无长度压缩）',
                '   - 维度压缩：D=256 → d=8',
                '   - 紧凑潜在表示 z'
            ]
        },
        {
            'title': '实验结果 - 结构重建',
            'content': [
                '数据集：AFDB-FS（588,318个结构，长度32-256）',
                '测试集：CASP14 和 CASP15',
                '',
                '主要结果：',
                '• RMSD 指标优于现有自编码器',
                '• 达到 SOTA 重建质量',
                '• 在潜在空间可实现准确的理化性质预测',
                '',
                '对比方法：',
                '• ESM3 structure autoencoder',
                '• VQ-VAE 方法',
                '• 其他离散编码方法'
            ]
        },
        {
            'title': '实验结果 - 蛋白质生成',
            'content': [
                'Protein Latent Diffusion Model (PLDM)：',
                '• 在 ProteinAE 学习的潜在空间上训练',
                '• 200M 参数，15层 DiT',
                '• 无需显式等变约束',
                '',
                '性能表现：',
                '• 显著优于先前的潜在空间方法',
                '• 性能匹敌经典的结构扩散模型（SDMs）',
                '• 缩小了潜在方法与结构方法之间的性能差距',
                '• 更高的设计性和多样性'
            ]
        },
        {
            'title': '与相关工作的对比',
            'content': [
                'vs. ESM3 Structure Autoencoder：',
                '• ESM3：SE(3) 空间，离散编码（VQ-VAE）',
                '• ProteinAE：E(3) 空间，连续编码',
                '',
                'vs. AlphaFold3：',
                '• 相似：非等变注意力架构',
                '• 不同：ProteinAE 专注于编码器-解码器',
                '',
                'vs. 其他扩散模型：',
                '• 在潜在空间操作，效率更高',
                '• 单一训练目标，简化优化'
            ]
        },
        {
            'title': '消融实验',
            'content': [
                '关键组件验证：',
                '• Bottleneck 维度影响（d=4, 8, 16）',
                '• 长度下采样比例（r=1, 2, 4）',
                '• LayerNorm vs KL 正则化',
                '• RoPE 位置编码的效果',
                '',
                '发现：',
                '• d=8 是性能和压缩的最佳平衡',
                '• LayerNorm 表现优于 KL 正则化',
                '• RoPE 对可变长度处理至关重要'
            ]
        },
        {
            'title': '结论与展望',
            'content': [
                '主要贡献：',
                '• 提出 ProteinAE，简化蛋白质结构编码',
                '• 连续潜在空间实现高效生成',
                '• SOTA 重建质量和生成性能',
                '',
                '未来方向：',
                '• 扩展到多链蛋白质和复合物',
                '• 结合序列信息进行联合建模',
                '• 应用于更多下游任务（柔性预测、功能预测等）',
                '• 探索更高效的潜在空间设计'
            ]
        }
    ]
    
    return outline

def fill_slides(prs, outline, paper_info):
    """填充幻灯片内容"""
    
    slide_idx = 1  # 从第2页开始（第1页是标题）
    
    for section in outline:
        if slide_idx >= len(prs.slides):
            break
        
        slide = prs.slides[slide_idx]
        
        # 查找标题占位符
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                # 如果是标题区域或占位符
                if shape.is_placeholder or len(text) < 50:
                    # 设置标题
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip():
                                run.text = section['title']
                                run.font.size = Pt(28)
                                run.font.bold = True
                                break
                else:
                    # 设置内容
                    content_text = '\n'.join(section['content'])
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = ""
                    
                    # 添加内容
                    p = shape.text_frame.paragraphs[0]
                    for line in section['content']:
                        run = p.add_run()
                        run.text = line + '\n'
                        run.font.size = Pt(16)
        
        slide_idx += 1
    
    return prs

if __name__ == "__main__":
    print("=" * 60)
    print("制作 ProteinAE 组会分享 PPT")
    print("=" * 60)
    
    # 创建PPT
    prs, paper_info = create_presentation()
    
    # 获取内容大纲
    outline = add_content_slides(prs, paper_info)
    
    # 填充内容
    prs = fill_slides(prs, outline, paper_info)
    
    # 保存
    output_path = '/data/ProteinAE_组会分享.pptx'
    prs.save(output_path)
    
    print(f"\n✅ PPT 已保存至: {output_path}")
    print(f"📊 共 {len(prs.slides)} 页")
    print("\n内容大纲：")
    print("1. 标题页")
    for i, section in enumerate(outline, 2):
        print(f"{i}. {section['title']}")
